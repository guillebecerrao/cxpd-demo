#!/usr/bin/env python3
"""E2E-SDPB — NexHealth / Coverwise — TW Executive Deck (English)"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.join(os.path.dirname(__file__), "E2E-SDPB-Coverwise-TW.pptx")

NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x1E, 0x88, 0xE5)
BLTL   = RGBColor(0xE3, 0xF2, 0xFD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFFW   = RGBColor(0xF4, 0xF6, 0xF8)
DARK   = RGBColor(0x21, 0x21, 0x21)
GRAY   = RGBColor(0x75, 0x75, 0x75)
LGRAY  = RGBColor(0xE0, 0xE0, 0xE0)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
LGRN   = RGBColor(0xE8, 0xF5, 0xE9)
AMBER  = RGBColor(0xE6, 0x5C, 0x00)
LAMB   = RGBColor(0xFF, 0xF8, 0xE1)
CYAN   = RGBColor(0x00, 0x83, 0x8E)
PURP   = RGBColor(0x4A, 0x14, 0x8C)
RED    = RGBColor(0xB7, 0x1C, 0x1C)
BLUELT = RGBColor(0x90, 0xCA, 0xF9)
BGRAY  = RGBColor(0xB0, 0xBE, 0xC5)
DARK2  = RGBColor(0x1A, 0x27, 0x3A)
DARK3  = RGBColor(0x0A, 0x14, 0x1F)
TEAL   = RGBColor(0x80, 0xCB, 0xC4)


def mk():
    p = Presentation()
    p.slide_width = Cm(33.87)
    p.slide_height = Cm(19.05)
    return p

def sl(p): return p.slides.add_slide(p.slide_layouts[6])

def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def box(s, l, t, w, h, fill):
    sh = s.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    return sh

def tx(s, text, l, t, w, h, sz=12, c=DARK, bold=False,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = s.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = c
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb

def mtx(s, lines, l, t, w, h, sz=12, c=DARK, bold=False,
        align=PP_ALIGN.LEFT, font="Calibri", sp=0):
    tb = s.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if sp and not first: p.space_before = Pt(sp)
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.color.rgb = c
        r.font.bold = bold; r.font.name = font
    return tb

def hdr(s, label):
    box(s, 0, 0, 33.87, 1.6, NAVY)
    tx(s, label, 1.5, 0.25, 30, 1.1, sz=19, c=WHITE, bold=True)


# ── SLIDE 1 — COVER ──────────────────────────────────────────────────────────
def s01(p):
    s = sl(p); bg(s, NAVY)
    box(s, 0, 0, 33.87, 0.3, BLUE)
    tx(s, "NexHealth  x  Coverwise", 1.5, 1.9, 25, 0.9, sz=12, c=BLUE)
    tx(s, "E2E-SDPB in Action", 1.5, 3.1, 28, 3.0, sz=50, c=WHITE, bold=True)
    tx(s, "From brief to live product in 5 days", 1.5, 7.0, 26, 1.4, sz=22, c=BLUELT)
    box(s, 1.5, 9.1, 8, 0.08, BLUE)
    tx(s, "Spec-Driven Product Building  |  End-to-End Framework", 1.5, 9.5, 30, 0.8, sz=11, c=GRAY)
    tx(s, "coverwise.becerra-ojeda.cl", 1.5, 10.5, 22, 0.8, sz=11, c=CYAN, font="Courier New")
    tx(s, "May 2026", 29.5, 17.8, 4, 0.8, sz=10, c=GRAY, align=PP_ALIGN.RIGHT)


# ── SLIDE 2 — THE BRIEF ──────────────────────────────────────────────────────
def s02(p):
    s = sl(p); bg(s, WHITE); hdr(s, "The Brief")
    for x, w, color, label, lines in [
        (1.5,  9.5, NAVY, "THE CLIENT",  ["NexHealth", "Digital health insurer", "LATAM operations", "500K active members", "Fast-growing young adult segment"]),
        (12.2, 12.0, RED,  "THE PROBLEM", ["Members can't tell if their", "procedure is covered when", "they need it.", "", "->  Call center overloaded", "->  NPS is low", "->  Members skip care"]),
        (25.4, 7.0,  CYAN, "THE USER",    ["25-40 years old", "Digital plan", "Mobile-first", "", '"Is this covered?"', "before going to the clinic"]),
    ]:
        box(s, x, 2.2, w, 7.0, OFFW)
        box(s, x, 2.2, w, 0.7, color)
        tx(s, label, x+0.35, 2.3, w-0.5, 0.55, sz=8.5, c=LGRAY, bold=True)
        mtx(s, lines, x+0.35, 3.15, w-0.5, 5.8, sz=12.5, c=DARK, sp=3)

    box(s, 0, 9.9, 33.87, 2.5, NAVY)
    tx(s, "THE PRODUCT", 1.8, 10.1, 8, 0.55, sz=8.5, c=BLUE, bold=True)
    tx(s, "Coverwise", 1.8, 10.75, 10, 0.9, sz=20, c=WHITE, bold=True)
    tx(s, "Conversational AI that answers coverage questions in natural language, in real time, from mobile.",
       13.0, 10.5, 19.5, 1.7, sz=13.5, c=BGRAY)


# ── SLIDE 3 — THE PROPOSAL ───────────────────────────────────────────────────
def s03(p):
    s = sl(p); bg(s, WHITE); hdr(s, "The Proposal: 5 Days, One Live Product")
    days = [
        ("DAY 1", "Shape",         "OST with\nprioritized\nopportunity",   NAVY),
        ("DAY 2", "Ideate",        "Ideas +\nAssumption\nTest Plan",        BLUE),
        ("DAY 3", "Validate",      "Evidence +\nChosen\nSolution",          CYAN),
        ("DAY 4", "Handoff",       "Specs +\nRelease\nPlan",                PURP),
        ("DAY 5", "Build\n& Deploy","Coverwise\nlive in\nproduction",       GREEN),
    ]
    cw, gap = 5.6, 0.5
    for i, (day, phase, out, col) in enumerate(days):
        x = 1.3 + i * (cw + gap)
        box(s, x, 2.0, cw, 10.5, OFFW)
        box(s, x, 2.0, cw, 1.4, col)
        tx(s, day, x+0.25, 2.1, cw-0.5, 0.8, sz=10, c=WHITE, bold=True)
        tx(s, phase, x+0.25, 3.8, cw-0.5, 2.2, sz=17, c=col, bold=True)
        tx(s, out, x+0.25, 6.3, cw-0.5, 4.0, sz=11.5, c=GRAY)
        if i < 4:
            tx(s, "->", x+cw+0.08, 5.5, 0.45, 1.0, sz=15, c=LGRAY, align=PP_ALIGN.CENTER)

    box(s, 0, 13.8, 33.87, 2.0, NAVY)
    tx(s, "coverwise.becerra-ojeda.cl  goes live on Day 5",
       1.8, 14.2, 30, 1.1, sz=15, c=WHITE, bold=True)


# ── SLIDE 4 — THE FRAMEWORK ──────────────────────────────────────────────────
def s04(p):
    s = sl(p); bg(s, NAVY)
    tx(s, "HOW E2E-SDPB WORKS", 1.5, 0.5, 25, 0.7, sz=10, c=BLUE, bold=True)
    tx(s, "The Framework", 1.5, 1.2, 28, 1.5, sz=30, c=WHITE, bold=True)
    phases = [
        ("Shape",           "Discover & prioritize\nopportunities",              "OST with\nprioritized leaf node",    BLUE),
        ("Ideate &\nValidate","Generate ideas,\ntest critical assumptions",       "Solution chosen\nwith evidence",     CYAN),
        ("Handoff to\nDelivery","Produce specs\ntraceable to evidence",           "Approved specs\n+ release plan",     PURP),
        ("Build",           "Claude Code builds\nfrom specs.\nEngineer reviews.", "Live product\nin production",       GREEN),
    ]
    bw, gap = 7.4, 0.25
    for i, (name, desc, out, col) in enumerate(phases):
        x = 0.8 + i * (bw + gap)
        box(s, x, 4.0, bw, 0.7, col)
        tx(s, name, x+0.2, 4.1, bw-0.4, 0.55, sz=12, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
        box(s, x, 4.7, bw, 5.0, DARK2)
        tx(s, desc, x+0.3, 4.95, bw-0.6, 2.8, sz=12, c=BGRAY)
        tx(s, "OUTPUT", x+0.3, 8.0, bw-0.6, 0.5, sz=8, c=col, bold=True)
        tx(s, out, x+0.3, 8.65, bw-0.6, 2.0, sz=11, c=WHITE)
        if i < 3:
            tx(s, "->", x+bw+0.03, 6.5, 0.3, 1.0, sz=17, c=BGRAY, align=PP_ALIGN.CENTER)
    box(s, 0, 12.5, 33.87, 2.0, DARK3)
    tx(s, '"Nothing gets built without a spec.  No spec without evidence.  Everything validated."',
       2.0, 12.85, 29.87, 1.3, sz=14, c=TEAL, italic=True, align=PP_ALIGN.CENTER)


# ── SLIDE 5 — DAY 1 + DAY 2 ─────────────────────────────────────────────────
def s05(p):
    s = sl(p); bg(s, WHITE); hdr(s, "Day 1 — Shape   +   Day 2 — Ideate")
    box(s, 1.5, 1.9, 0.4, 14.2, NAVY)
    tx(s, "SHAPE", 2.15, 2.1, 5, 0.6, sz=10, c=NAVY, bold=True)
    mtx(s, [
        "Kickoff with NexHealth: mission, users, constraints, API availability",
        "Claude (PM agent) populates corporate context live during the session",
        "Opportunity mapping workshop: 5 opportunities identified",
        "Claude builds the full OST during the workshop in real time",
        "Sizing & prioritization: leaf node selected with explicit criteria",
    ], 2.15, 2.9, 14.5, 6.0, sz=11.5, c=DARK, sp=4)
    box(s, 2.15, 9.1, 14.6, 2.8, BLTL)
    box(s, 2.15, 9.1, 0.27, 2.8, BLUE)
    tx(s, "Leaf node elected:", 2.65, 9.3, 14, 0.55, sz=8.5, c=BLUE, bold=True)
    tx(s, '"The member doesn\'t know if their procedure is covered when they need it"',
       2.65, 10.0, 13.8, 1.7, sz=12, c=NAVY, italic=True)
    tx(s, "KEY ARTIFACTS", 2.15, 12.2, 14, 0.5, sz=8.5, c=GRAY, bold=True)
    mtx(s, ["contexto/estrategia/ost.md  (v1)",
            "shape/analisis-oportunidades.md",
            "shape/oportunidades-sizing-priorizacion.md"],
        2.15, 12.85, 14.5, 2.5, sz=10, c=GRAY, font="Courier New", sp=3)

    box(s, 17.5, 1.9, 0.4, 14.2, BLUE)
    tx(s, "IDEATE", 18.15, 2.1, 5, 0.6, sz=10, c=BLUE, bold=True)
    mtx(s, [
        "Ideation: HMW, analogies, assumption inversion",
        "5 solution ideas evaluated on impact / effort / technical risk",
        "Conversational AI assistant selected for deep-dive validation",
        "Critical assumptions mapped by risk level",
        "Test plan: 5 user interviews + Wizard of Oz session",
    ], 18.15, 2.9, 14.5, 6.0, sz=11.5, c=DARK, sp=4)
    box(s, 18.15, 9.1, 14.6, 2.8, LGRN)
    box(s, 18.15, 9.1, 0.27, 2.8, GREEN)
    tx(s, "Critical assumptions to test:", 18.65, 9.3, 14, 0.55, sz=8.5, c=GREEN, bold=True)
    mtx(s, ["Members trust AI for coverage decisions  [HIGH risk]",
            "Natural language > forms  [HIGH risk]",
            "Coverage API has enough granularity  [MEDIUM risk]"],
        18.65, 10.0, 13.8, 1.7, sz=12, c=DARK, sp=3)
    tx(s, "KEY ARTIFACTS", 18.15, 12.2, 14, 0.5, sz=8.5, c=GRAY, bold=True)
    mtx(s, ["ideate-validate/ideas.md",
            "ideate-validate/solution-assumptions.md",
            "ideate-validate/test-plan.md"],
        18.15, 12.85, 14.5, 2.5, sz=10, c=GRAY, font="Courier New", sp=3)

    box(s, 0, 16.3, 33.87, 0.75, OFFW)
    tx(s, "Agents:  PM agent  |  OST Facilitator skill  |  Designer agent",
       1.5, 16.42, 30, 0.55, sz=10, c=GRAY)


# ── SLIDE 6 — VALIDATE ───────────────────────────────────────────────────────
def s06(p):
    s = sl(p); bg(s, WHITE); hdr(s, "Day 3 — Validate: Evidence-Based Decision")
    box(s, 1.5, 1.95, 30.9, 1.5, NAVY)
    tx(s, "Method: 5 user interviews (members, 25-40 yo, digital plans) + Wizard of Oz test",
       2.0, 2.2, 29, 1.0, sz=13, c=WHITE)
    box(s, 1.5, 3.85, 30.9, 0.75, DARK)
    for lbl, x, w in [("Assumption", 2.0, 13.5), ("Result", 16.5, 5.5), ("Finding", 23.0, 9.0)]:
        tx(s, lbl, x, 3.95, w, 0.6, sz=9.5, c=WHITE, bold=True)
    rows = [
        ("Members trust AI for coverage decisions",    "CONFIRMED", '"If it says yes, I go." Conditional trust is acceptable.',              LGRN, GREEN),
        ("Prefer natural language over forms",         "CONFIRMED", "Forms cause drop-off at field 3. Natural input feels like texting.",      LGRN, GREEN),
        ("Use it from mobile at moment of care",       "CONFIRMED", '4/5 participants: "before going to the clinic" as the key trigger.',     LGRN, GREEN),
        ("Coverage API has enough granularity",        "PARTIAL",   "API returns by FONASA code, not name. Requires translation layer (+1d).", LAMB, AMBER),
    ]
    for i, (assumption, result, finding, bgc, ac) in enumerate(rows):
        y = 4.6 + i * 2.1
        box(s, 1.5, y, 30.9, 2.1, OFFW if i % 2 == 0 else WHITE)
        box(s, 1.5, y, 0.2, 2.1, ac)
        tx(s, assumption, 2.0, y+0.1, 13.5, 1.9, sz=11, c=DARK)
        tx(s, result, 16.5, y+0.6, 6.5, 0.85, sz=11, c=ac, bold=True)
        tx(s, finding, 23.0, y+0.1, 9.2, 1.9, sz=10, c=DARK)
    box(s, 0, 13.05, 33.87, 2.2, NAVY)
    tx(s, "CHOSEN SOLUTION:", 1.8, 13.25, 10, 0.55, sz=8.5, c=BLUE, bold=True)
    tx(s, "Coverwise MVP -- Conversational AI + Procedure Name Resolver (FONASA translation layer)",
       1.8, 13.9, 31, 1.0, sz=14, c=WHITE, bold=True)
    box(s, 0, 15.25, 33.87, 0.7, OFFW)
    tx(s, "Evidence: EV-001 to EV-005  |  Agents: PM agent  |  Designer agent  |  Interview Snapshot skill",
       1.5, 15.35, 30, 0.55, sz=10, c=GRAY)


# ── SLIDE 7 — HANDOFF ────────────────────────────────────────────────────────
def s07(p):
    s = sl(p); bg(s, WHITE); hdr(s, "Day 4 — Handoff to Delivery: Specs That Build Themselves")
    tx(s, "Claude (SDD agent) decomposes the solution into specs traceable to the OST. Engineer approves. PM validates intent.",
       1.5, 1.9, 30.9, 0.8, sz=11.5, c=GRAY)
    box(s, 1.5, 3.0, 30.9, 0.7, DARK)
    for lbl, x, w in [("Feature", 2.0, 5), ("P", 7.5, 1.5), ("Name", 9.5, 8), ("Acceptance Criteria (summary)", 18.5, 13.5)]:
        tx(s, lbl, x, 3.1, w, 0.55, sz=9.5, c=WHITE, bold=True)
    feats = [
        ("FEAT-001", "P0", "Coverage query interface",
         "Free-text input, autocomplete, mobile-first (touch >= 44px). Max 5 suggestions.", RED),
        ("FEAT-002", "P0", "Coverage result display",
         '"Covered / Not Covered / Partial" + copay limits and conditions in plain language.', RED),
        ("FEAT-003", "P0", "Procedure name resolver",
         "NL -> FONASA code translation. Handles typos, variants, accents. Mock data for demo.", RED),
        ("FEAT-004", "P1", "Fallback to human agent",
         "When resolver confidence < threshold: escalate to live agent with full query context.", BLUE),
    ]
    for i, (feat, pri, name, ac, pc) in enumerate(feats):
        y = 3.7 + i * 1.9
        box(s, 1.5, y, 30.9, 1.9, OFFW if i % 2 == 0 else WHITE)
        tx(s, feat, 2.0, y+0.5, 5, 0.7, sz=10.5, c=GRAY, font="Courier New")
        box(s, 7.55, y+0.55, 1.2, 0.7, pc)
        tx(s, pri, 7.6, y+0.6, 1.1, 0.6, sz=9, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tx(s, name, 9.5, y+0.3, 8.5, 1.3, sz=12, c=DARK, bold=True)
        tx(s, ac, 18.5, y+0.15, 13.5, 1.65, sz=10, c=GRAY)
    box(s, 1.5, 11.35, 30.9, 4.35, DARK2)
    tx(s, "FEAT-001 -- sample spec excerpt", 2.0, 11.52, 22, 0.55, sz=8.5, c=BLUE, bold=True)
    mtx(s, [
        "user story:   As a NexHealth member, I want to type my procedure in plain language",
        "              to get an immediate, clear coverage answer.",
        "",
        "acceptance:   Text input activates after 3 chars  |  max 5 suggestions shown",
        '              If no match: "Procedure not found. Speak to an agent?"',
        "              Touch target >= 44px  |  Handles: typos, abbreviations, accents",
    ], 2.0, 12.2, 30.0, 3.3, sz=10, c=TEAL, font="Courier New", sp=2)
    box(s, 0, 15.7, 33.87, 0.65, OFFW)
    tx(s, "Every spec traces to the OST  |  Agents: SDD agent  |  Engineer agent",
       1.5, 15.78, 30, 0.5, sz=10, c=GRAY)


# ── SLIDE 8 — BUILD & DEPLOY ─────────────────────────────────────────────────
def s08(p):
    s = sl(p); bg(s, NAVY)
    tx(s, "DAY 5  |  BUILD & DEPLOY", 1.5, 0.4, 25, 0.75, sz=10, c=BLUE, bold=True)
    tx(s, "The Engineer Doesn't Write a Line", 1.5, 1.15, 30, 1.55, sz=28, c=WHITE, bold=True)
    tx(s, "STACK", 1.5, 3.15, 6, 0.55, sz=8.5, c=BLUE, bold=True)
    stack = [
        ("Frontend",  "React + Tailwind CSS  (mobile-first)"),
        ("Backend",   "Python / FastAPI"),
        ("AI layer",  "Claude API  (claude-sonnet-4-6)"),
        ("Coverage",  "Mock FONASA dataset  (JSON, demo)"),
        ("Infra",     "Google Cloud Run"),
        ("Domain",    "coverwise.becerra-ojeda.cl"),
    ]
    for i, (layer, tech) in enumerate(stack):
        y = 3.85 + i * 1.3
        box(s, 1.5, y, 3.5, 1.1, DARK2)
        tx(s, layer, 1.75, y+0.2, 3.0, 0.7, sz=9, c=BLUE, bold=True)
        tx(s, tech, 5.3, y+0.25, 12.5, 0.7, sz=11, c=WHITE)
    box(s, 19.0, 2.5, 14.0, 11.2, DARK3)
    tx(s, "HOW IT HAPPENS", 19.5, 2.7, 13, 0.55, sz=8.5, c=BLUE, bold=True)
    for i, step in enumerate([
        "Claude Code reads approved specs from the repo",
        "Builds: FEAT-003 -> FEAT-001 -> FEAT-002 -> FEAT-004",
        "Engineer reviews each output -- no code from scratch",
        "gcloud run deploy  (Cloud Run, us-central1)",
        "DNS pointed: coverwise.becerra-ojeda.cl is live",
        "Demo to NexHealth: repo walkthrough + live URL",
    ]):
        tx(s, f"{i+1}.  {step}", 19.8, 3.5+i*1.55, 12.8, 1.4, sz=11, c=BGRAY)
    box(s, 0, 14.3, 33.87, 2.5, GREEN)
    tx(s, "LIVE IN PRODUCTION", 1.8, 14.55, 18, 0.65, sz=10.5, c=WHITE, bold=True)
    tx(s, "coverwise.becerra-ojeda.cl", 1.8, 15.2, 32, 1.1, sz=22, c=WHITE, bold=True, font="Courier New")


# ── SLIDE 9 — THE REPO ───────────────────────────────────────────────────────
def s09(p):
    s = sl(p); bg(s, WHITE); hdr(s, "The SDPB-Context Repo: Product Memory in Git")
    tx(s, "REPO STRUCTURE", 1.5, 2.0, 12, 0.55, sz=8.5, c=NAVY, bold=True)
    for i, (line, col, bold) in enumerate([
        ("cxpd-demo/",                                     NAVY,  True),
        ("  +-- CLAUDE.md                <- agent instructions",   BLUE,  False),
        ("  +-- framework.md             <- E2E-SDPB rules",       BLUE,  False),
        ("  +-- contexto/",                                DARK,  True),
        ("  |   +-- corporativo/         <- kickoff context",      GRAY,  False),
        ("  |   +-- estrategia/ost.md    <- single source of truth",GRAY, False),
        ("  |   +-- evidencia/           <- EV-001 to EV-005",     GRAY,  False),
        ("  +-- shape/                   <- Day 1 artifacts",      DARK,  True),
        ("  +-- ideate-validate/         <- Days 2-3 artifacts",   DARK,  True),
        ("  +-- handoff-delivery/        <- Day 4 specs + plan",   DARK,  True),
        ("  +-- bitacora/                <- full cycle log",       GRAY,  False),
    ]):
        tx(s, line, 1.5, 2.75+i*1.05, 17, 1.0, sz=10, c=col, bold=bold, font="Courier New")
    for i, (title, desc, col) in enumerate([
        ("No re-explaining",  "Every session starts with full context. The agent reads the repo before acting.", BLUE),
        ("No lost decisions", "Every choice -- from opportunity to spec -- is documented and traceable.",        CYAN),
        ("Full traceability", "Spec -> solution -> assumption test -> evidence -> OST. End to end.",             GREEN),
        ("Async-native",      "PM, Designer, Engineer, and agents work async. The repo is the sync point.",      PURP),
    ]):
        y = 2.0 + i * 3.2
        box(s, 19.2, y, 14.0, 3.0, OFFW)
        box(s, 19.2, y, 0.27, 3.0, col)
        tx(s, title, 19.8, y+0.3, 13.2, 0.75, sz=13, c=col, bold=True)
        tx(s, desc,  19.8, y+1.2, 13.2, 1.65, sz=11, c=DARK)
    box(s, 0, 14.65, 33.87, 0.75, OFFW)
    tx(s, '"The repo is not a document archive. It\'s the operating system of the product."',
       1.8, 14.78, 30, 0.55, sz=11, c=GRAY, italic=True)


# ── SLIDE 10 — HOW WE WORK TOGETHER ─────────────────────────────────────────
def s10(p):
    s = sl(p); bg(s, NAVY)
    tx(s, "THE PROPOSAL", 1.5, 0.5, 15, 0.65, sz=10, c=BLUE, bold=True)
    tx(s, "How We Work Together", 1.5, 1.2, 28, 1.55, sz=28, c=WHITE, bold=True)
    tx(s, "THE TRIAD", 1.5, 3.2, 10, 0.55, sz=8.5, c=BLUE, bold=True)
    for i, (role, desc, col) in enumerate([
        ("Product Manager", "Facilitates discovery. Owns the OST and product brief.\nCoordinates with client stakeholders.", BLUE),
        ("Product Designer", "Leads UX research, ideation sessions,\nand prototype validation.", CYAN),
        ("Product Engineer", "Technical viability, spec approval,\nbuild supervision with Claude Code.", GREEN),
    ]):
        y = 3.95 + i * 2.6
        box(s, 1.5, y, 0.38, 2.1, col)
        tx(s, role, 2.2, y+0.1, 14, 0.75, sz=13, c=col, bold=True)
        tx(s, desc, 2.2, y+0.95, 14, 1.5, sz=11, c=BGRAY)
    box(s, 18.0, 2.5, 14.5, 12.7, DARK3)
    tx(s, "OPERATING MODEL", 18.5, 2.7, 13.5, 0.55, sz=8.5, c=BLUE, bold=True)
    for i, pt in enumerate([
        "Client codebase stays in client repos",
        "SDPB-Context repo shared with client team",
        "AI agents augment the triad -- not replace it",
        "Every artifact committed to git, full history",
        "AI-provider agnostic (Claude, GPT, Gemini)",
        "5-day cycle is repeatable per product increment",
    ]):
        box(s, 18.7, 3.7+i*1.75, 0.45, 0.45, BLUE)
        tx(s, pt, 19.5, 3.62+i*1.75, 12.5, 0.95, sz=12, c=WHITE)
    box(s, 0, 15.25, 33.87, 1.85, BLUE)
    tx(s, "Ready to run a 5-day cycle with a real client engagement?",
       1.8, 15.55, 25, 0.95, sz=16, c=WHITE, bold=True)
    tx(s, "Guillermo Becerra  |  E2E-SDPB",
       26.5, 16.15, 7, 0.7, sz=10, c=BLTL, align=PP_ALIGN.RIGHT)


# ── MAIN ─────────────────────────────────────────────────────────────────────
def main():
    deck = mk()
    s01(deck); s02(deck); s03(deck); s04(deck); s05(deck)
    s06(deck); s07(deck); s08(deck); s09(deck); s10(deck)
    deck.save(OUT)
    print(f"Saved -> {OUT}")

if __name__ == "__main__":
    main()
